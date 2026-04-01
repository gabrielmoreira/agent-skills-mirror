#!/usr/bin/env python
"""Build a PPTX from a JSON/YAML slide spec."""

from __future__ import annotations

import argparse
import json
import os
import sys
from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches, Pt


def _load_spec(path: str) -> Dict[str, Any]:
    """Internal helper to load spec.
    
    Args:
        path (str): Filesystem path used by this operation.
    
    Returns:
        Dict[str, Any]: Result produced by this function.
    
    Raises:
        RuntimeError: Raised when an execution error occurs.
    
    Note:
        This is a private helper used internally by the module/class.
    """
    _, ext = os.path.splitext(path.lower())
    with open(path, "r", encoding="utf-8") as f:
        raw = f.read()

    if ext in {".yaml", ".yml"}:
        try:
            import yaml  # type: ignore
        except Exception as exc:  # pragma: no cover
            raise RuntimeError(
                "YAML spec requires PyYAML. Install with: python -m pip install pyyaml"
            ) from exc
        return yaml.safe_load(raw)

    return json.loads(raw)


def _set_widescreen(prs: Presentation) -> None:
    """Internal helper to set widescreen.
    
    Args:
        prs (Presentation): Input value for prs.
    
    Returns:
        None: This method does not return a value.
    
    Note:
        This is a private helper used internally by the module/class.
    """
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def _add_title_slide(prs: Presentation, slide_spec: Dict[str, Any]) -> None:
    """Internal helper to add title slide.
    
    Args:
        prs (Presentation): Input value for prs.
        slide_spec (Dict[str, Any]): Input value for slide spec.
    
    Returns:
        None: This method does not return a value.
    
    Note:
        This is a private helper used internally by the module/class.
    """
    layout = prs.slide_layouts[0] if prs.slide_layouts else None
    slide = prs.slides.add_slide(layout) if layout else prs.slides.add_slide(prs.slide_layouts[5])

    title = slide_spec.get("title", "")
    subtitle = slide_spec.get("subtitle", "")

    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
        tx.text_frame.text = title

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    elif subtitle:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), prs.slide_width - Inches(1), Inches(1))
        tx.text_frame.text = subtitle


def _add_bullets_slide(prs: Presentation, slide_spec: Dict[str, Any]) -> None:
    """Internal helper to add bullets slide.
    
    Args:
        prs (Presentation): Input value for prs.
        slide_spec (Dict[str, Any]): Input value for slide spec.
    
    Returns:
        None: This method does not return a value.
    
    Note:
        This is a private helper used internally by the module/class.
    """
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)

    title = slide_spec.get("title", "")
    if slide.shapes.title:
        slide.shapes.title.text = title
    elif title:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.7))
        tx.text_frame.text = title

    body = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body = shape
            break

    if body is None:
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), prs.slide_width - Inches(1.6), prs.slide_height - Inches(2))

    tf = body.text_frame
    tf.clear()

    bullets = slide_spec.get("bullets", [])
    for item in bullets:
        if isinstance(item, dict):
            text = item.get("text", "")
            level = int(item.get("level", 0))
        else:
            text = str(item)
            level = 0
        p = tf.add_paragraph()
        p.text = text
        p.level = max(level, 0)


def _add_two_col_slide(prs: Presentation, slide_spec: Dict[str, Any]) -> None:
    """Internal helper to add two col slide.
    
    Args:
        prs (Presentation): Input value for prs.
        slide_spec (Dict[str, Any]): Input value for slide spec.
    
    Returns:
        None: This method does not return a value.
    
    Note:
        This is a private helper used internally by the module/class.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide_spec.get("title", "")
    if title:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.7))
        tx.text_frame.text = title

    left = slide_spec.get("left", [])
    right = slide_spec.get("right", [])

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), (prs.slide_width / 2) - Inches(0.75), prs.slide_height - Inches(2))
    right_box = slide.shapes.add_textbox((prs.slide_width / 2) + Inches(0.25), Inches(1.3), (prs.slide_width / 2) - Inches(0.75), prs.slide_height - Inches(2))

    for items, box in [(left, left_box), (right, right_box)]:
        tf = box.text_frame
        tf.clear()
        for item in items:
            p = tf.add_paragraph()
            p.text = str(item)
            p.level = 0


def _add_image_slide(prs: Presentation, slide_spec: Dict[str, Any]) -> None:
    """Internal helper to add image slide.
    
    Args:
        prs (Presentation): Input value for prs.
        slide_spec (Dict[str, Any]): Input value for slide spec.
    
    Returns:
        None: This method does not return a value.
    
    Raises:
        ValueError: Raised when an execution error occurs.
    
    Note:
        This is a private helper used internally by the module/class.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    image_path = slide_spec.get("image")
    if not image_path:
        raise ValueError("Image slide requires 'image' path")

    slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    title = slide_spec.get("title", "")
    caption = slide_spec.get("caption", "")
    if title:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.6))
        tx.text_frame.text = title
        tx.text_frame.paragraphs[0].font.size = Pt(28)
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.8), prs.slide_width - Inches(1), Inches(0.6))
        tx.text_frame.text = caption
        tx.text_frame.paragraphs[0].font.size = Pt(16)


def build_pptx(spec: Dict[str, Any], out_path: str, template: str | None, widescreen: bool) -> None:
    """Build pptx.
    
    Args:
        spec (Dict[str, Any]): Input value for spec.
        out_path (str): Filesystem path for out path.
        template (str | None): Input value for template.
        widescreen (bool): Input value for widescreen.
    
    Returns:
        None: This method does not return a value.
    
    Raises:
        ValueError: Raised when an execution error occurs.
    """
    prs = Presentation(template) if template else Presentation()
    if widescreen:
        _set_widescreen(prs)

    slides = spec.get("slides", [])
    for slide_spec in slides:
        stype = slide_spec.get("type", "bullets")
        if stype == "title":
            _add_title_slide(prs, slide_spec)
        elif stype == "bullets":
            _add_bullets_slide(prs, slide_spec)
        elif stype == "two_col":
            _add_two_col_slide(prs, slide_spec)
        elif stype == "image":
            _add_image_slide(prs, slide_spec)
        else:
            raise ValueError(f"Unknown slide type: {stype}")

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)


def main() -> int:
    """Main.
    
    Args:
        None.
    
    Returns:
        int: Result produced by this function.
    """
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True, help="Path to JSON/YAML spec")
    parser.add_argument("--out", required=True, help="Output .pptx path")
    parser.add_argument("--template", help="Optional .pptx template")
    parser.add_argument("--widescreen", action="store_true", help="Use 16:9 slide size")

    args = parser.parse_args()

    spec = _load_spec(args.spec)
    build_pptx(spec, args.out, args.template, args.widescreen)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
